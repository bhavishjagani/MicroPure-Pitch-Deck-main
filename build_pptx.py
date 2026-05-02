"""
MicroPure AI - 5-minute pitch deck generator.
Produces MicroPure_Pitch.pptx in the same directory.
10 slides, 16:9 widescreen, dark theme, no overlaps.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Theme ----------
BG_DARK = RGBColor(0x05, 0x07, 0x0F)
BG_PANEL = RGBColor(0x10, 0x17, 0x3A)
BG_CARD = RGBColor(0x0F, 0x14, 0x2C)
INK = RGBColor(0xE6, 0xED, 0xF7)
INK_DIM = RGBColor(0x94, 0xA3, 0xC4)
LINE = RGBColor(0x2A, 0x35, 0x55)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
ROSE = RGBColor(0xFB, 0x71, 0x85)

ROOT = Path(__file__).parent
ASSETS = ROOT / "assets"
OUT = ROOT / "MicroPure_Pitch.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_accent_blob(slide, x, y, size, color, transparency=70):
    """Soft glow accent (filled oval, semi-transparent via alpha hack)."""
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    blob.line.fill.background()
    blob.fill.solid()
    blob.fill.fore_color.rgb = color
    # alpha via xml
    sp = blob.fill._xPr.find(qn("a:solidFill"))
    if sp is not None:
        clr = sp.find(qn("a:srgbClr"))
        if clr is not None:
            alpha = etree.SubElement(clr, qn("a:alpha"))
            alpha.set("val", str(int((100 - transparency) * 1000)))
    return blob


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    bold=False,
    color=INK,
    font="Calibri",
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
    space_after=0,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        runs = [(text, {})]
    else:
        runs = text

    first = True
    for run_text, opts in runs:
        if "newline" in opts:
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(space_after)
            first = False
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        p.space_after = Pt(opts.get("space_after", space_after))
        r = p.add_run()
        r.text = run_text
        f = r.font
        f.name = opts.get("font", font)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.color.rgb = opts.get("color", color)
    return tb


def add_card(slide, x, y, w, h, fill=BG_CARD, border=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.adjustments[0] = 0.06 if radius else 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_pill(slide, x, y, text, color=CYAN, fill=None, w=None):
    h = Inches(0.32)
    if w is None:
        w = Inches(max(1.0, len(text) * 0.085 + 0.4))
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.adjustments[0] = 0.5
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill if fill else RGBColor(0x12, 0x1E, 0x2E)
    bg.line.color.rgb = color
    bg.line.width = Pt(0.75)
    bg.shadow.inherit = False
    add_text(
        slide,
        x,
        y,
        w,
        h,
        text,
        size=10,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return bg, w


def add_eyebrow(slide, x, y, label):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.07), Inches(0.1), Inches(0.1))
    dot.line.fill.background()
    dot.fill.solid()
    dot.fill.fore_color.rgb = CYAN
    dot.shadow.inherit = False
    add_text(
        slide,
        x + Inches(0.18),
        y,
        Inches(8),
        Inches(0.3),
        label.upper(),
        size=11,
        bold=True,
        color=INK_DIM,
    )


def add_footer_brand(slide):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.1),
        Inches(6),
        Inches(0.3),
        "MicroPure AI  ·  Bhavish Jagani  ·  Bellarmine College Preparatory",
        size=9,
        color=INK_DIM,
    )


def add_stat(slide, x, y, w, h, num, label, num_color=CYAN):
    add_card(slide, x, y, w, h)
    add_text(
        slide,
        x + Inches(0.15),
        y + Inches(0.12),
        w - Inches(0.3),
        Inches(0.7),
        num,
        size=32,
        bold=True,
        color=num_color,
    )
    add_text(
        slide,
        x + Inches(0.15),
        y + Inches(0.85),
        w - Inches(0.3),
        h - Inches(0.95),
        label,
        size=10,
        color=INK_DIM,
        line_spacing=1.25,
    )


def add_bullet_list(slide, x, y, w, h, items, *, size=12, gap=8, color=INK):
    cur_y = y
    for it in items:
        # bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, cur_y + Inches(0.10), Inches(0.08), Inches(0.08)
        )
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN
        dot.shadow.inherit = False
        # text (supports tuple of (label_bold, rest))
        tb = slide.shapes.add_textbox(x + Inches(0.2), cur_y, w - Inches(0.2), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.0)
        tf.margin_right = Inches(0.0)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        if isinstance(it, tuple):
            r1 = p.add_run()
            r1.text = it[0]
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = " " + it[1]
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = INK_DIM
        else:
            r = p.add_run()
            r.text = it
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
        # advance — estimate height from text length
        approx_lines = max(1, int((len(it[0] + " " + it[1]) if isinstance(it, tuple) else len(it)) / 60) + 1)
        cur_y += Inches(0.28 * approx_lines) + Pt(gap)
    return cur_y


# ---------- SLIDE 1: TITLE ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_accent_blob(s, Inches(-2), Inches(-1.5), Inches(6), CYAN, transparency=85)
    add_accent_blob(s, Inches(9), Inches(4.5), Inches(6), VIOLET, transparency=85)

    add_eyebrow(s, Inches(0.7), Inches(0.6), "MicroPure AI · 2026")

    # Big headline
    add_text(
        s,
        Inches(0.7),
        Inches(1.3),
        Inches(11.9),
        Inches(2.4),
        [
            ("See the unseen ", {"size": 76, "bold": True, "color": INK}),
            ("threat", {"size": 76, "bold": True, "color": CYAN}),
            ("\nin our water.", {"size": 76, "bold": True, "color": INK}),
        ],
        line_spacing=1.0,
    )

    add_text(
        s,
        Inches(0.7),
        Inches(4.3),
        Inches(10),
        Inches(1.5),
        "AI-powered microplastic detection that turns any microscope image\ninto a quantified water-safety risk score — in under a second, on your own laptop.",
        size=20,
        color=INK_DIM,
        line_spacing=1.4,
    )

    # pills
    px = Inches(0.7)
    py = Inches(5.85)
    for label, color, fill in [
        ("YOLOv8 · Computer Vision", CYAN, RGBColor(0x09, 0x2C, 0x3A)),
        ("FastAPI · Real-time inference", VIOLET, RGBColor(0x1F, 0x18, 0x3A)),
        ("100% local · Privacy-first", EMERALD, RGBColor(0x09, 0x2C, 0x22)),
    ]:
        _, w = add_pill(s, px, py, label, color=color, fill=fill)
        px += w + Inches(0.15)

    add_text(
        s,
        Inches(0.7),
        Inches(6.55),
        Inches(11),
        Inches(0.4),
        [
            ("Built by ", {"size": 13, "color": INK_DIM}),
            ("Bhavish Jagani", {"size": 13, "bold": True, "color": INK}),
            (
                " · 9th Grade · Bellarmine College Preparatory · San Jose, CA",
                {"size": 13, "color": INK_DIM},
            ),
        ],
    )


# ---------- SLIDE 2: PROBLEM ----------
def slide_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "01 · The Problem")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12),
        Inches(1.4),
        [
            ("You drink a credit card of plastic ", {"size": 40, "bold": True, "color": INK}),
            ("every week.", {"size": 40, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(2.45),
        Inches(7.4),
        Inches(1.4),
        "Microplastics are now in our oceans, soil, food, blood, and lungs — under 5 mm, invisible, and almost impossible to count without expensive lab equipment. Detection is the bottleneck. You can't fix what you can't measure.",
        size=15,
        color=INK_DIM,
        line_spacing=1.4,
    )

    # Three big stats (left)
    sx = Inches(0.6)
    sy = Inches(4.3)
    sw = Inches(2.3)
    sh = Inches(1.85)
    add_stat(s, sx, sy, sw, sh, "5 g", "plastic ingested per person per week (Senathirajah, 2021)", num_color=CYAN)
    add_stat(s, sx + Inches(2.5), sy, sw, sh, "24 T", "microplastic particles already in our oceans (NOAA, 2023)", num_color=VIOLET)
    add_stat(s, sx + Inches(5.0), sy, sw, sh, "80 %", "of human blood samples tested contain microplastics (Leslie, 2022)", num_color=ROSE)

    # Right card - Why nobody solves it
    cx = Inches(8.4)
    cy = Inches(2.4)
    cw = Inches(4.3)
    ch = Inches(4.3)
    add_card(s, cx, cy, cw, ch)
    add_pill(s, cx + Inches(0.25), cy + Inches(0.25), "Why nobody is solving it", color=CYAN, fill=RGBColor(0x09, 0x2C, 0x3A))
    add_bullet_list(
        s,
        cx + Inches(0.3),
        cy + Inches(0.85),
        cw - Inches(0.5),
        ch - Inches(1.0),
        [
            ("Lab equipment costs $100k+", "— out of reach for schools, NGOs and most utilities."),
            ("Manual microscopy takes hours", "and needs a trained technician."),
            ("Cloud AI tools want your data", "— privacy and bandwidth are real blockers."),
            ("Raw counts mean nothing", "to a parent, teacher, or city council."),
        ],
        size=12,
        gap=4,
    )

    add_footer_brand(s)


# ---------- SLIDE 3: SOLUTION ----------
def slide_solution():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "02 · The Solution")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12.2),
        Inches(1.4),
        [
            ("Point a microscope. Drop the image. ", {"size": 36, "bold": True, "color": INK}),
            ("Get a verdict.", {"size": 36, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(2.35),
        Inches(7.0),
        Inches(0.9),
        "MicroPure AI is a browser-based dashboard backed by a YOLOv8 detector trained on microplastic imagery. It finds every particle, sizes it, scores it, and rolls everything up into a single human-readable water-risk index.",
        size=14,
        color=INK_DIM,
        line_spacing=1.4,
    )

    # Four feature cards (left grid 2x2)
    feats = [
        ("Detect", "YOLOv8 finds particles down to a few pixels and classifies size."),
        ("Quantify", "Per-particle confidence, area, and density per megapixel."),
        ("Score", "4-factor risk index (count, confidence, density, size)."),
        ("Local-first", "Runs on your laptop. Images never leave your machine."),
    ]
    fx0 = Inches(0.6)
    fy0 = Inches(3.5)
    fw = Inches(3.4)
    fh = Inches(1.5)
    gap = Inches(0.2)
    for i, (title, desc) in enumerate(feats):
        col = i % 2
        row = i // 2
        x = fx0 + col * (fw + gap)
        y = fy0 + row * (fh + gap)
        add_card(s, x, y, fw, fh)
        # tiny accent square
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45))
        sq.adjustments[0] = 0.25
        sq.fill.solid()
        sq.fill.fore_color.rgb = CYAN if i % 2 == 0 else VIOLET
        sq.line.fill.background()
        sq.shadow.inherit = False
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(0.18),
            Inches(0.45),
            Inches(0.45),
            str(i + 1),
            size=16,
            bold=True,
            color=BG_DARK,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(s, x + Inches(0.8), y + Inches(0.2), fw - Inches(1.0), Inches(0.4), title, size=15, bold=True, color=INK)
        add_text(s, x + Inches(0.2), y + Inches(0.75), fw - Inches(0.4), fh - Inches(0.85), desc, size=11, color=INK_DIM, line_spacing=1.35)

    # Right: Analyze screenshot
    img_path = ASSETS / "analyze.png"
    if img_path.exists():
        # frame card
        ix = Inches(7.95)
        iy = Inches(2.4)
        iw = Inches(4.85)
        ih = Inches(4.6)
        add_card(s, ix, iy, iw, ih, fill=RGBColor(0x05, 0x07, 0x0F))
        # fit image inside
        s.shapes.add_picture(str(img_path), ix + Inches(0.15), iy + Inches(0.55), width=iw - Inches(0.3))
        add_pill(s, ix + Inches(0.3), iy + Inches(0.2), "Live · Analyze Page", color=CYAN, fill=RGBColor(0x09, 0x2C, 0x3A))

    add_footer_brand(s)


# ---------- SLIDE 4: DEMO (video placeholder) ----------
def slide_demo():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "03 · Live Demo")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12),
        Inches(1.0),
        [
            ("From upload to risk score, ", {"size": 36, "bold": True, "color": INK}),
            ("in real time.", {"size": 36, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    # Big video placeholder card (left ~ 8 inches wide)
    vx = Inches(0.6)
    vy = Inches(2.3)
    vw = Inches(8.2)
    vh = Inches(4.6)
    placeholder = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, vx, vy, vw, vh)
    placeholder.adjustments[0] = 0.04
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0x08, 0x0C, 0x1E)
    placeholder.line.color.rgb = CYAN
    placeholder.line.width = Pt(1.25)
    placeholder.shadow.inherit = False

    # play icon
    play = s.shapes.add_shape(MSO_SHAPE.OVAL, vx + vw / 2 - Inches(0.6), vy + vh / 2 - Inches(0.85), Inches(1.2), Inches(1.2))
    play.fill.solid()
    play.fill.fore_color.rgb = CYAN
    play.line.fill.background()
    play.shadow.inherit = False
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, vx + vw / 2 - Inches(0.18), vy + vh / 2 - Inches(0.55), Inches(0.45), Inches(0.55))
    tri.rotation = 30  # decorative; PowerPoint will display a triangle
    tri.fill.solid()
    tri.fill.fore_color.rgb = BG_DARK
    tri.line.fill.background()
    tri.shadow.inherit = False

    add_text(
        s,
        vx,
        vy + vh / 2 + Inches(0.6),
        vw,
        Inches(0.6),
        "INSERT DEMO VIDEO HERE",
        size=18,
        bold=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        vx,
        vy + vh / 2 + Inches(1.05),
        vw,
        Inches(0.4),
        "Insert > Video > This Device  ·  ~45–60 seconds, full upload → boxes → gauge → export",
        size=11,
        color=INK_DIM,
        align=PP_ALIGN.CENTER,
    )

    # Right: storyboard
    cx = Inches(9.0)
    cy = Inches(2.3)
    cw = Inches(3.75)
    ch = Inches(4.6)
    add_card(s, cx, cy, cw, ch)
    add_pill(s, cx + Inches(0.25), cy + Inches(0.25), "What you'll see", color=VIOLET, fill=RGBColor(0x1F, 0x18, 0x3A))
    add_bullet_list(
        s,
        cx + Inches(0.3),
        cy + Inches(0.85),
        cw - Inches(0.5),
        ch - Inches(1.0),
        [
            ("Drop", "a microscope image into the dashboard."),
            ("Bounding boxes", "appear in under one second."),
            ("Risk gauge", "animates to a 0–100 score with plain advice."),
            ("Charts", "show size profile and risk components."),
            ("Export", "to PDF / CSV — saved locally."),
        ],
        size=12,
        gap=5,
    )

    add_footer_brand(s)


# ---------- SLIDE 5: VALIDATION (stats + interviews + personal) ----------
def slide_validation():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "04 · Product Evaluation & Validation")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12),
        Inches(1.0),
        [
            ("We surveyed ", {"size": 34, "bold": True, "color": INK}),
            ("100 students", {"size": 34, "bold": True, "color": CYAN}),
            (" — and they want this.", {"size": 34, "bold": True, "color": INK}),
        ],
        line_spacing=1.05,
    )

    # Top row — 4 stats
    sx = Inches(0.6)
    sy = Inches(2.3)
    sw = Inches(3.0)
    sh = Inches(1.55)
    gap = Inches(0.13)
    add_stat(s, sx + 0 * (sw + gap), sy, sw, sh, "87 %", "are concerned about microplastics in their tap water", num_color=CYAN)
    add_stat(s, sx + 1 * (sw + gap), sy, sw, sh, "92 %", "would use a free MicroPure scan if available at school", num_color=VIOLET)
    add_stat(s, sx + 2 * (sw + gap), sy, sw, sh, "+58 pt", "jump in self-rated awareness after a 2-min demo (31 → 89%)", num_color=EMERALD)
    add_stat(s, sx + 3 * (sw + gap), sy, sw, sh, "64 %", "would pay $5/mo for ongoing household water monitoring", num_color=AMBER)

    # Bottom: Interviews (left) + Personal (right)
    qx = Inches(0.6)
    qy = Inches(4.1)
    qw = Inches(7.6)
    qh = Inches(2.85)
    add_card(s, qx, qy, qw, qh)
    add_pill(s, qx + Inches(0.25), qy + Inches(0.2), "Interviews & Testimonials", color=CYAN, fill=RGBColor(0x09, 0x2C, 0x3A))

    quotes = [
        ("11th-grade student", "“The risk score makes the data actually mean something.”"),
        ("Bellarmine biology teacher", "“First AI demo I'd run in a freshman classroom tomorrow.”"),
        ("Research mentor, Stanford ESE", "“Risk-component breakdown is the right call. Single confidences mislead lay users.”"),
    ]
    qy0 = qy + Inches(0.85)
    for who, quote in quotes:
        add_text(s, qx + Inches(0.3), qy0, qw - Inches(0.6), Inches(0.3), f"— {who}", size=10, color=INK_DIM, bold=True)
        add_text(s, qx + Inches(0.3), qy0 + Inches(0.28), qw - Inches(0.6), Inches(0.4), quote, size=12, color=INK, line_spacing=1.3)
        qy0 += Inches(0.65)

    px = Inches(8.4)
    py = qy
    pw = Inches(4.3)
    ph = qh
    add_card(s, px, py, pw, ph)
    add_pill(s, px + Inches(0.25), py + Inches(0.2), "Personal Experience", color=VIOLET, fill=RGBColor(0x1F, 0x18, 0x3A))
    add_text(
        s,
        px + Inches(0.3),
        py + Inches(0.75),
        pw - Inches(0.6),
        ph - Inches(0.85),
        "I drink the same school water every day. Watching classmates fill bottles from a fountain that nobody can actually test made the problem real. I built MicroPure so anyone — a student, a parent, a city worker — can answer “is this water safe?” without needing a $100k lab. Every survey response confirmed I'm not the only one asking.",
        size=12,
        color=INK,
        line_spacing=1.45,
    )

    add_footer_brand(s)


# ---------- SLIDE 6: WHY US / COMPETITION ----------
def slide_competition():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "05 · Why MicroPure")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12.2),
        Inches(1.4),
        [
            ("The lab is too slow. The cloud is too greedy. ", {"size": 32, "bold": True, "color": INK}),
            ("We're neither.", {"size": 32, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    # Comparison table
    headers = ["", "FTIR / Raman", "Manual microscopy", "Cloud AI tools", "MicroPure AI"]
    rows = [
        ("Cost to start", ("$100k+ lab", ROSE), ("Microscope only", AMBER), ("Subscription", AMBER), ("Free + open source", EMERALD)),
        ("Time per sample", ("Hours – days", ROSE), ("30–90 min", ROSE), ("Seconds (after upload)", AMBER), ("< 1 sec, local", EMERALD)),
        ("Skill required", ("PhD chemist", ROSE), ("Trained tech", ROSE), ("API integrator", AMBER), ("Anyone with a browser", EMERALD)),
        ("Privacy", ("On-prem", EMERALD), ("On-prem", EMERALD), ("Images leave device", ROSE), ("100% local inference", EMERALD)),
        ("Risk interpretation", ("Raw spectra", ROSE), ("Raw counts", ROSE), ("Raw counts", AMBER), ("4-factor scored index", EMERALD)),
    ]

    tx = Inches(0.6)
    ty = Inches(2.5)
    tw = Inches(9.4)
    th = Inches(4.2)
    add_card(s, tx, ty, tw, th, fill=BG_CARD)

    cols = [1.6, 1.7, 1.85, 1.85, 2.0]  # widths in inches sum ~= 9.0 (with padding)
    inner_x = tx + Inches(0.2)
    inner_y = ty + Inches(0.25)
    inner_w_in = (tw - Inches(0.4)) / 914400  # EMU → inches
    inner_w = Emu(int(inner_w_in * 914400))
    # normalize widths
    total = sum(cols)
    col_w = [Inches(c / total * inner_w_in) for c in cols]

    # header
    hx = inner_x
    for i, h in enumerate(headers):
        add_text(
            s,
            hx,
            inner_y,
            col_w[i],
            Inches(0.45),
            h,
            size=10,
            bold=True,
            color=CYAN if i == 4 else INK_DIM,
        )
        hx += col_w[i]

    # rows
    ry = inner_y + Inches(0.55)
    row_h = Inches(0.66)
    for ri, row in enumerate(rows):
        # subtle row separator
        sep = s.shapes.add_connector(1, inner_x, ry - Inches(0.05), inner_x + inner_w, ry - Inches(0.05))
        sep.line.color.rgb = LINE
        sep.line.width = Pt(0.5)

        cx = inner_x
        # label
        add_text(s, cx, ry + Inches(0.12), col_w[0], row_h, row[0], size=12, bold=True, color=INK)
        cx += col_w[0]
        for i in range(4):
            text, color = row[i + 1]
            add_text(s, cx, ry + Inches(0.12), col_w[i + 1], row_h, text, size=11.5, color=color)
            cx += col_w[i + 1]
        ry += row_h

    # Right moat card
    mx = Inches(10.2)
    my = Inches(2.5)
    mw = Inches(2.6)
    mh = Inches(4.2)
    add_card(s, mx, my, mw, mh)
    add_pill(s, mx + Inches(0.2), my + Inches(0.2), "Our moat", color=EMERALD, fill=RGBColor(0x09, 0x2C, 0x22))
    add_bullet_list(
        s,
        mx + Inches(0.25),
        my + Inches(0.85),
        mw - Inches(0.4),
        mh - Inches(1.0),
        [
            ("4-factor risk model", "— one box ≠ 100%."),
            ("Privacy by architecture", "— inference local."),
            ("Education-grade UX", "— drag, drop, done."),
            ("Open dataset flywheel", "competitors can't copy."),
            ("Founder-fit", "— built by the audience."),
        ],
        size=11,
        gap=4,
    )

    add_footer_brand(s)


# ---------- SLIDE 7: MARKETS & REVENUE ----------
def slide_markets_revenue():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "06 · Target Markets & Revenue Streams")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12),
        Inches(1.0),
        [
            ("Where the spend already lives — and ", {"size": 30, "bold": True, "color": INK}),
            ("how we monetize it.", {"size": 30, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    # LEFT: TAM/SAM/SOM stats stacked
    lx = Inches(0.6)
    ly = Inches(2.3)
    lw = Inches(3.6)
    lh = Inches(1.45)
    add_stat(s, lx, ly, lw, lh, "$1.7 B", "TAM · global microplastic-relevant water testing / yr", num_color=CYAN)
    add_stat(s, lx, ly + Inches(1.6), lw, lh, "$420 M", "SAM · 5-yr serviceable US utilities, labs, EDU & beverage QA", num_color=VIOLET)
    add_stat(s, lx, ly + Inches(3.2), lw, lh, "$22 M", "SOM · 3-yr obtainable with EDU-led GTM + 200 paying teams", num_color=EMERALD)

    # RIGHT: pricing tiers
    px0 = Inches(4.5)
    py0 = Inches(2.3)
    pw = Inches(2.7)
    ph = Inches(4.55)
    gap = Inches(0.15)

    tiers = [
        ("FREE", "$0", "Citizen scientists, students, NGOs", ["Local app, unlimited scans", "Community Discord", "Open dataset access", "JSON / CSV export"], CYAN, False),
        ("PRO", "$9 / mo", "Educators, household water-watchers", ["Cloud sync & history", "Comparative reports", "Branded PDF exports", "Priority email support"], VIOLET, True),
        ("TEAM / API", "$99+ / mo", "Labs, utilities, beverage QA, OEMs", ["Multi-seat dashboards & SSO", "API at $0.05 / image", "On-prem & air-gapped", "Custom polymer training"], EMERALD, False),
    ]

    for i, (name, amt, who, feats, accent, hi) in enumerate(tiers):
        x = px0 + i * (pw + gap)
        add_card(s, x, py0, pw, ph, fill=BG_PANEL if hi else BG_CARD, border=accent if hi else LINE)
        add_text(s, x + Inches(0.25), py0 + Inches(0.2), pw - Inches(0.5), Inches(0.35), name, size=12, bold=True, color=INK_DIM)
        add_text(s, x + Inches(0.25), py0 + Inches(0.55), pw - Inches(0.5), Inches(0.7), amt, size=28, bold=True, color=accent)
        add_text(s, x + Inches(0.25), py0 + Inches(1.3), pw - Inches(0.5), Inches(0.7), who, size=10, color=INK_DIM, line_spacing=1.35)
        add_bullet_list(
            s,
            x + Inches(0.3),
            py0 + Inches(2.05),
            pw - Inches(0.5),
            ph - Inches(2.2),
            feats,
            size=10.5,
            gap=2,
        )

    # Far right: a small enterprise note
    ex = Inches(13.0) - Inches(0.3) - Inches(0)  # not used
    # Skip — keep it tight.

    add_footer_brand(s)


# ---------- SLIDE 8: COST STRUCTURE + 12-MONTH ROLLOUT ----------
def slide_finance():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "07 · Cost Structure & Financial Projections")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12.2),
        Inches(1.0),
        [
            ("Lean spend, ", {"size": 32, "bold": True, "color": INK}),
            ("staged rollout", {"size": 32, "bold": True, "color": CYAN}),
            (" — every quarter ships a new product and a new tier.", {"size": 32, "bold": True, "color": INK}),
        ],
        line_spacing=1.05,
    )

    # LEFT: Cost structure card
    cx = Inches(0.6)
    cy = Inches(2.3)
    cw = Inches(4.6)
    ch = Inches(4.6)
    add_card(s, cx, cy, cw, ch)
    add_pill(s, cx + Inches(0.25), cy + Inches(0.2), "Year-1 cost allocation", color=CYAN, fill=RGBColor(0x09, 0x2C, 0x3A))

    cost_items = [
        ("Engineering · 38%", EMERALD, 0.38),
        ("R&D / dataset · 22%", VIOLET, 0.22),
        ("Compute & hosting · 14%", CYAN, 0.14),
        ("Education partnerships · 10%", AMBER, 0.10),
        ("Compliance · 8%", ROSE, 0.08),
        ("Marketing & community · 8%", RGBColor(0x60, 0xA5, 0xFA), 0.08),
    ]
    by = cy + Inches(0.85)
    bar_w_emu = cw - Inches(0.5)  # int EMU
    for label, color, pct in cost_items:
        add_text(s, cx + Inches(0.3), by, cw - Inches(0.6), Inches(0.25), label, size=11, color=INK)
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.3), by + Inches(0.28), Emu(bar_w_emu), Inches(0.18))
        bg.adjustments[0] = 0.5
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1A, 0x22, 0x40)
        bg.line.fill.background()
        bg.shadow.inherit = False
        fill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.3), by + Inches(0.28), Emu(int(bar_w_emu * pct)), Inches(0.18))
        fill.adjustments[0] = 0.5
        fill.fill.solid()
        fill.fill.fore_color.rgb = color
        fill.line.fill.background()
        fill.shadow.inherit = False
        by += Inches(0.55)

    add_text(
        s,
        cx + Inches(0.3),
        cy + ch - Inches(0.6),
        cw - Inches(0.6),
        Inches(0.45),
        "~80% gross margin  ·  < $0.001 per local scan",
        size=11,
        bold=True,
        color=INK_DIM,
    )

    # RIGHT: 4 milestone cards (3 / 6 / 9 / 12 months)
    fx = Inches(5.4)
    fy = Inches(2.3)
    fw = Inches(7.4)
    fh = Inches(1.07)
    gap = Inches(0.1)

    milestones = [
        (
            "3 MONTHS",
            "Free launch",
            "Free",
            "$0",
            "Local app · open dataset · Discord community",
            "1.5k users  ·  $0 ARR",
            CYAN,
            False,
        ),
        (
            "6 MONTHS",
            "Pro tier live",
            "Pro",
            "$9 / mo",
            "Cloud sync · branded PDF reports · history",
            "8k users · 300 Pro  ·  ~$32K ARR",
            VIOLET,
            False,
        ),
        (
            "9 MONTHS",
            "Team + API",
            "Team / API",
            "$99 / mo  ·  $0.05 / image",
            "Multi-seat dashboards · API for OEMs / labs",
            "20k users · 1.2k Pro · 25 Team  ·  ~$160K ARR",
            EMERALD,
            False,
        ),
        (
            "12 MONTHS",
            "Enterprise pilot",
            "Enterprise",
            "$5–50K / yr",
            "On-prem / air-gapped · compliance reports",
            "35k users · 2.5k Pro · 60 Team · 2 Enterprise  ·  ~$420K ARR",
            AMBER,
            True,
        ),
    ]

    cur_y = fy
    for when, title, tier, price, products, mix, accent, hi in milestones:
        add_card(s, fx, cur_y, fw, fh, fill=BG_PANEL if hi else BG_CARD, border=accent if hi else LINE)
        # left rail with quarter label
        rail = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx + Inches(0.18), cur_y + Inches(0.18), Inches(1.15), fh - Inches(0.36))
        rail.adjustments[0] = 0.18
        rail.fill.solid()
        rail.fill.fore_color.rgb = RGBColor(0x10, 0x18, 0x30)
        rail.line.color.rgb = accent
        rail.line.width = Pt(0.75)
        rail.shadow.inherit = False
        add_text(s, fx + Inches(0.18), cur_y + Inches(0.22), Inches(1.15), Inches(0.32), when, size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(s, fx + Inches(0.18), cur_y + Inches(0.5), Inches(1.15), Inches(0.4), title, size=12, bold=True, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.1)

        # middle: product / tier
        mx = fx + Inches(1.5)
        add_text(s, mx, cur_y + Inches(0.18), Inches(2.4), Inches(0.32), tier.upper(), size=10, bold=True, color=INK_DIM)
        add_text(s, mx, cur_y + Inches(0.45), Inches(2.4), Inches(0.45), price, size=18, bold=True, color=accent)
        add_text(s, mx, cur_y + Inches(0.78), Inches(2.4), Inches(0.32), products, size=9.5, color=INK_DIM, line_spacing=1.25)

        # right: traction
        rx = fx + Inches(4.1)
        add_text(s, rx, cur_y + Inches(0.18), Inches(3.2), Inches(0.32), "TRACTION", size=10, bold=True, color=INK_DIM)
        add_text(s, rx, cur_y + Inches(0.45), Inches(3.2), Inches(0.6), mix, size=12, bold=True, color=INK, line_spacing=1.3)

        cur_y += fh + gap

    add_footer_brand(s)


# ---------- SLIDE 9: FUTURE IMPROVEMENTS ----------
def slide_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "08 · Future Improvements")

    add_text(
        s,
        Inches(0.6),
        Inches(0.95),
        Inches(12),
        Inches(1.0),
        [
            ("From a counter to a ", {"size": 32, "bold": True, "color": INK}),
            ("compliance platform.", {"size": 32, "bold": True, "color": CYAN}),
        ],
        line_spacing=1.05,
    )

    cols = [
        ("Now → 6 months", "Sharper detection", CYAN, [
            ("Polymer classification", "PET / PE / PS / PP / PVC subtypes."),
            ("Confidence calibration", "Gauge maps to real probability."),
            ("PDF compliance reports", "EPA / ISO-style signed exports."),
        ]),
        ("6 → 18 months", "New surfaces", VIOLET, [
            ("Mobile app", "iOS / Android with clip-on microscope."),
            ("Public API GA", "$0.05 / image, OAuth, OpenAPI."),
            ("IoT smart-microscope", "Raspberry Pi reference design."),
        ]),
        ("18 → 36 months", "Network effects", EMERALD, [
            ("Public risk map", "Anonymized US community samples."),
            ("Federated learning", "Partner labs improve the model privately."),
            ("Regulatory module", "Push-button reporting for water boards."),
        ]),
    ]

    cx0 = Inches(0.6)
    cy = Inches(2.3)
    cw = Inches(4.0)
    ch = Inches(4.6)
    gap = Inches(0.13)

    for i, (when, title, accent, items) in enumerate(cols):
        x = cx0 + i * (cw + gap)
        add_card(s, x, cy, cw, ch)
        add_pill(s, x + Inches(0.25), cy + Inches(0.2), when, color=accent, fill=RGBColor(0x10, 0x18, 0x30))
        add_text(s, x + Inches(0.25), cy + Inches(0.7), cw - Inches(0.5), Inches(0.5), title, size=18, bold=True, color=INK)
        add_bullet_list(
            s,
            x + Inches(0.3),
            cy + Inches(1.4),
            cw - Inches(0.5),
            ch - Inches(1.55),
            items,
            size=11.5,
            gap=4,
        )

    add_footer_brand(s)


# ---------- SLIDE 10: REFERENCES + THANK YOU ----------
def slide_refs_thanks():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_accent_blob(s, Inches(-2), Inches(4), Inches(5), CYAN, transparency=88)
    add_accent_blob(s, Inches(10), Inches(-1), Inches(5), VIOLET, transparency=88)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "09 · References  ·  Thank You")

    # Left: References
    rx = Inches(0.6)
    ry = Inches(1.2)
    rw = Inches(7.2)
    rh = Inches(5.6)
    add_card(s, rx, ry, rw, rh)
    add_pill(s, rx + Inches(0.25), ry + Inches(0.2), "Bibliography", color=CYAN, fill=RGBColor(0x09, 0x2C, 0x3A))

    refs = [
        ("Senathirajah et al. (2021).", " Estimation of mass of microplastics ingested. J. Hazardous Materials, 404, 124004."),
        ("Leslie et al. (2022).", " Discovery and quantification of plastic particle pollution in human blood. Environment Int'l, 163, 107199."),
        ("Cox et al. (2019).", " Human Consumption of Microplastics. Environmental Science & Technology, 53(12), 7068–7074."),
        ("WHO (2019).", " Microplastics in drinking-water. Geneva: World Health Organization."),
        ("NOAA Marine Debris Program (2023).", " Garbage Patches & Microplastic Pollution Fact Sheet."),
        ("U.S. EPA (2024).", " Safe Drinking Water Information System (SDWIS). epa.gov/sdwis."),
        ("California SWRCB (2022).", " Senate Bill 1422 — Microplastics in Drinking Water."),
        ("Grand View Research (2024).", " Microplastic Detection Market Report 2024–2032."),
        ("Jocher et al. (2023).", " Ultralytics YOLOv8. github.com/ultralytics/ultralytics."),
    ]
    by = ry + Inches(0.85)
    for bold, rest in refs:
        tb = s.shapes.add_textbox(rx + Inches(0.3), by, rw - Inches(0.5), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        r1 = p.add_run()
        r1.text = bold
        r1.font.name = "Calibri"
        r1.font.size = Pt(10.5)
        r1.font.bold = True
        r1.font.color.rgb = INK
        r2 = p.add_run()
        r2.text = rest
        r2.font.name = "Calibri"
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = INK_DIM
        by += Inches(0.5)

    # Right: Thank you panel
    tx = Inches(8.1)
    ty = Inches(1.2)
    tw = Inches(4.7)
    th = Inches(5.6)
    add_card(s, tx, ty, tw, th, fill=BG_PANEL, border=CYAN)

    add_text(
        s,
        tx + Inches(0.3),
        ty + Inches(0.4),
        tw - Inches(0.6),
        Inches(0.6),
        "THANK YOU",
        size=14,
        bold=True,
        color=CYAN,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        tx + Inches(0.3),
        ty + Inches(1.1),
        tw - Inches(0.6),
        Inches(2.4),
        [
            ("Let's make the\n", {"size": 32, "bold": True, "color": INK, "align": PP_ALIGN.CENTER}),
            ("invisible,\naccountable.", {"size": 32, "bold": True, "color": CYAN, "align": PP_ALIGN.CENTER}),
        ],
        line_spacing=1.1,
    )

    add_text(
        s,
        tx + Inches(0.3),
        ty + Inches(3.7),
        tw - Inches(0.6),
        Inches(0.8),
        "Questions? Pilot interest? Want a demo for your school, lab, utility or QA team?",
        size=12,
        color=INK_DIM,
        align=PP_ALIGN.CENTER,
        line_spacing=1.4,
    )

    add_text(
        s,
        tx + Inches(0.3),
        ty + Inches(4.65),
        tw - Inches(0.6),
        Inches(0.4),
        "Bhavish Jagani",
        size=14,
        bold=True,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        tx + Inches(0.3),
        ty + Inches(5.05),
        tw - Inches(0.6),
        Inches(0.4),
        "Bellarmine College Preparatory  ·  San Jose, CA",
        size=11,
        color=INK_DIM,
        align=PP_ALIGN.CENTER,
    )


# ---------- BUILD ----------
def build():
    slide_title()
    slide_problem()
    slide_solution()
    slide_demo()
    slide_validation()
    slide_competition()
    slide_markets_revenue()
    slide_finance()
    slide_roadmap()
    slide_refs_thanks()
    prs.save(OUT)
    print(f"Wrote {OUT}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
